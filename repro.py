import pptx
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

presentation = pptx.Presentation("template.pptx")

chart_data = ChartData()
chart_data.categories = ['Yes', 'No']
chart_data.add_series('Series 1', (42, 24))

for chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT]:
    for has_data_labels in [True, False]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        placeholder = slide.placeholders[13]
        chart = placeholder.insert_chart(chart_type, chart_data).chart

        for plot in chart.plots:
            plot.has_data_labels = has_data_labels

presentation.save("output.pptx")
